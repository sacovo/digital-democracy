"""
Misc. functions that are used by the project

These functions provide functionallity for the application but are not directly linked
to either models or views.
"""
from dataclasses import dataclass
import csv
import io
import re
import secrets
from itertools import zip_longest
from typing import List

from bs4 import BeautifulSoup
from django.conf import settings
from django.contrib.auth import get_user_model
from django.core.validators import validate_email
from django.utils.translation import gettext as _
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Cm, Pt


@dataclass
class Change:
    """
    This class represents a single change to a paper. It is a simple data store class
    that is used by the methods below to apply changes to a paper.
    """

    start: int
    end: int
    content: str  # eg: <ins>...</ins> or <del>...</del>


def add_classes_to_tags(tag_name, class_list, html) -> str:
    """
    Adds the classes in class_list to every element in the html with the given tag_name.
    Returns the html with the modifications as a new string.
    """
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup.find_all(tag_name):
        tag["class"] = tag.get("class", []) + class_list
    return str(soup)


def add_lite_classes(text) -> str:
    """
    Add the ice-del/ice-ins and ice-cts classes to the del and ins tags in the
    given html. This needs to be used in order to render the changes with colors
    and get the functionallity of ice.
    """
    text = add_classes_to_tags("del", ["ice-del", "ice-cts"], text)
    text = add_classes_to_tags("ins", ["ice-ins", "ice-cts"], text)
    return text


def apply_change(original_text: str, change: Change) -> str:
    """
    Apply the given change to the original text and return the result.
    """
    return "".join(
        (original_text[: change.start], change.content, original_text[change.end :])
    )


def update_change(change: Change, applied: Change):
    """
    Update the text positions in the first change after the second change has been applied.
    """
    if change.start < applied.start:
        return
    delta = len(applied.content) - (applied.end - applied.start)
    change.start += delta
    change.end += delta


def create_changes_of_amendment(text: str) -> List[Change]:
    """text contains changes marked up with <del> or <ins> tags.
    Subtraction of 11, because <del></del> is 11 characters.
    """
    changes = []
    parser = re.compile(r"(\<del>(.*?)\<\/del>)|(\<ins>(.*?)\<\/ins>)")

    matches = parser.finditer(text)
    counter = 0
    for match in matches:
        content = match.group()
        if content.startswith("<del>"):
            changes.append(
                Change(match.start() - counter, match.end() - 11 - counter, content)
            )
            counter += len("<del></del>")
        else:
            changes.append(
                Change(match.start() - counter, match.start() - counter, content)
            )
            counter += len(content)

    return changes


def create_modified_text(original_text: str, amendments) -> str:
    """
    Generate an html with del and ins tags for every amendment in amendments.
    """
    changes = []
    for amendment in amendments:
        changes += create_changes_of_amendment(amendment.content)
    modified_text = original_text

    changes.sort(key=lambda change: change.start)
    changes.reverse()
    for change in changes:
        modified_text = apply_change(modified_text, change)
    return modified_text


def index_of_first_change(content):
    """
    Returns the index of the first occurrence of a change
    """
    first_insert = content.find("<ins")
    first_deletion = content.find("<del")

    if first_insert == -1:
        return first_deletion

    if first_deletion == -1:
        return first_insert

    return min(first_deletion, first_insert)


def index_of_last_change(content):
    """
    Returns the index where the last change ends.
    """
    return max(content.rfind("</ins>"), content.rfind("</del>"))


def extract_content(content):
    """
    Takes a string with markup for changes (<ins> and <del> tags) and returns
    a string that contains these tags plus some context
    """
    start_index = index_of_first_change(content)

    end_index = index_of_last_change(content) + 6

    content_before = content[0:start_index]

    sentences = re.split(r"(\.|\?|\!)", content_before)

    sentences = list(zip_longest(sentences[::2], sentences[1::2], fillvalue=""))

    sentence_before = ""

    for sentence in reversed(sentences[-3:]):
        sentence_before = sentence[0] + sentence[1] + sentence_before

    content_after = content[end_index:]

    sentences = re.split(r"(\.|\?|\!)", content_after)

    sentences = list(zip_longest(sentences[::2], sentences[1::2], fillvalue=""))

    sentence_after = ""

    for sentence in sentences[:3]:
        sentence_after = sentence_after + sentence[0] + sentence[1]

    return sentence_before + content[start_index:end_index] + sentence_after


def import_users_from_csv(csv_file):
    """
    Import the given csv file into the database
    """
    csv_file = io.TextIOWrapper(csv_file)
    csv_reader = csv.DictReader(csv_file)

    imported_users = [get_user_model()(**row) for row in csv_reader]

    for new_user in imported_users:
        validate_email(new_user.email)

    for new_user in imported_users:
        password = secrets.token_urlsafe(17)
        new_user.set_password(password)
        new_user.save()

        new_user.email_user(
            _("New digital-democracy account"),
            settings.NEW_USER_MAIL.format(user=new_user, password=password),
            fail_silently=True,
        )

    return imported_users


def generate_powerpoint(paper):
    """
    Generates a pp-presentation with all amendments of the given paper.
    """
    # pylint: disable=R0914, R0915, E1101

    # Create new presentation
    prs = Presentation()

    # Create layouts
    t_layout = prs.slide_layouts[6]
    a_layout = prs.slide_layouts[1]

    # Add the title slide to the prs
    title_slide = prs.slides.add_slide(t_layout)
    titles_txtbox = title_slide.shapes.add_textbox(
        Cm(2), Cm(10), Cm(21), Cm(5)
    )  # (left, top, width, height)
    titles_txtframe = titles_txtbox.text_frame
    titles_txtframe.word_wrap = True
    p0_t_txtfrm = titles_txtframe.paragraphs[0]
    p0_t_txtfrm.text = "\n".join(
        (translation.title for translation in paper.translation_set.all())
    )
    p0_t_txtfrm.font.size = Pt(22)
    p0_t_txtfrm.font.bold = True
    p0_t_txtfrm.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p0_t_txtfrm.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p0_t_txtfrm.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Underline title
    underline = title_slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Cm(23), Cm(9.5), Cm(2.25), Cm(9.5)
    )
    line = LineFormat(underline)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Set background
    background = title_slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 0, 0)

    # Set title slide
    title_txtbx = title_slide.shapes.add_textbox(Cm(2), Cm(6.5), Cm(21), Cm(3))
    subtitle_tf = title_txtbx.text_frame
    subtitle = subtitle_tf.paragraphs[0]
    subtitle.text = "PAPER TITLE:"
    subtitle.font.size = Pt(60)
    subtitle.font.bold = True
    subtitle.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Generating amendment slides
    for i, amendment in enumerate(
        paper.amendment_set.filter(
            language_code=paper.translation_set.first().language_code
        )
    ):
        amendment_slide = prs.slides.add_slide(a_layout)

        amendment_nr_textbox = amendment_slide.shapes.add_textbox(
            Cm(0.25), Cm(0), Cm(4), Cm(2.8)
        )
        paper_title_textbox = amendment_slide.shapes.add_textbox(
            Cm(4.5), Cm(0.2), Cm(20), Cm(3)
        )

        txz_frame = amendment_nr_textbox.text_frame
        p0_amendment_nr = txz_frame.paragraphs[0]
        p0_amendment_nr.text = f"A{i + 1}:"
        p0_amendment_nr.font.size = Pt(70)
        p0_amendment_nr.font.bold = True
        p0_amendment_nr.font.color.rgb = RGBColor(0x0, 0x0, 0x0)

        # Paper title on amendment page
        title_textframe = paper_title_textbox.text_frame
        p0_amendment_paper_title = title_textframe.paragraphs[0]

        p0_amendment_paper_title.text = "\n".join(
            (translation.title for translation in paper.translation_set.all())
        )
        p0_amendment_paper_title.font.size = Pt(18)
        p0_amendment_paper_title.font.bold = True
        p0_amendment_paper_title.font.color.rgb = RGBColor(0x0, 0x0, 0x0)
        p0_amendment_paper_title.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        p0_amendment_paper_title.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # Underline title
        u_line2 = amendment_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, Cm(25), Cm(3), Cm(0.5), Cm(3)
        )
        u_line2 = LineFormat(u_line2)
        u_line2.fill.solid()
        u_line2.fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Removing unused placeholder
        txtbx = amendment_slide.shapes[0]
        txtbx.element.getparent().remove(txtbx.element)

        # Adding body bullet points
        body = amendment_slide.shapes.placeholders[1]
        body.text_frame.text = amendment.title + "\n"

        for translation in amendment.translation_list():
            bullet = body.text_frame.add_paragraph()
            bullet.text = translation.title + "\n"

    return prs
